from __future__ import annotations

import glob
import hashlib
import json
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from .config import AppConfig
from .io_utils import ensure_workspace, now_stamp, read_text_file, safe_slug, write_text_file
from .llm import LLMClient

TEXT_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".rst",
    ".yaml",
    ".yml",
    ".csv",
}
JSON_EXTENSIONS = {".json"}
HTML_EXTENSIONS = {".html", ".htm"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
LEGACY_PPT_EXTENSIONS = {".ppt"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg"}
VIDEO_EXTENSIONS = {".mp4", ".mov", ".mkv", ".webm", ".avi"}

ALL_SUPPORTED_EXTENSIONS = (
    TEXT_EXTENSIONS
    | JSON_EXTENSIONS
    | HTML_EXTENSIONS
    | PDF_EXTENSIONS
    | DOCX_EXTENSIONS
    | PPTX_EXTENSIONS
    | LEGACY_PPT_EXTENSIONS
    | IMAGE_EXTENSIONS
    | AUDIO_EXTENSIONS
    | VIDEO_EXTENSIONS
)


@dataclass(frozen=True)
class IngestItem:
    source_path: Path
    output_path: Path
    source_type: str
    extractor: str


@dataclass(frozen=True)
class IngestSummary:
    ingested: list[IngestItem]
    skipped_existing: list[Path]
    skipped_unsupported: list[Path]
    errors: list[str]


def normalize_source_input(value: str) -> str:
    item = str(value or "").strip()
    if len(item) >= 2 and item[0] == item[-1] and item[0] in {'"', "'"}:
        item = item[1:-1].strip()
    return item


def normalize_source_inputs(inputs: list[str]) -> list[str]:
    return [item for item in (normalize_source_input(value) for value in inputs) if item]


def _expand_sources(inputs: list[str]) -> list[Path]:
    files: list[Path] = []
    for item in normalize_source_inputs(inputs):
        has_glob = any(token in item for token in ("*", "?", "["))
        if has_glob:
            for match in glob.glob(item, recursive=True):
                path = Path(match)
                if path.is_file():
                    files.append(path.resolve())
            continue

        path = Path(item)
        if path.is_file():
            files.append(path.resolve())
            continue
        if path.is_dir():
            files.extend(candidate.resolve() for candidate in path.rglob("*") if candidate.is_file())
            continue

    return sorted(set(files))


def resolve_sources(inputs: list[str], *, max_files: int | None = None) -> list[Path]:
    files = _expand_sources(inputs)
    if max_files is not None:
        files = files[:max_files]
    return files


def _target_path(raw_dir: Path, source_path: Path) -> Path:
    digest = hashlib.sha1(str(source_path).encode("utf-8")).hexdigest()[:10]
    slug = safe_slug(source_path.stem, max_len=60)
    return raw_dir / "ingested" / f"{slug}-{digest}.md"


def _extract_text_like(source_path: Path) -> tuple[str, str]:
    return read_text_file(source_path), "plain_text"


def _extract_json(source_path: Path) -> tuple[str, str]:
    raw = read_text_file(source_path).lstrip("\ufeff")
    try:
        parsed = json.loads(raw)
        pretty = json.dumps(parsed, ensure_ascii=False, indent=2, sort_keys=True)
        return pretty, "json_pretty"
    except json.JSONDecodeError:
        return raw, "json_raw_fallback"


def _extract_html(source_path: Path) -> tuple[str, str]:
    try:
        from markdownify import markdownify
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency: markdownify. Run `pip install markdownify beautifulsoup4`."
        ) from exc
    html = read_text_file(source_path)
    markdown = markdownify(html, heading_style="ATX", strip=["script", "style"])
    return markdown.strip() or html, "html_to_markdown"


def _require_llm(llm: LLMClient | None, *, reason: str) -> LLMClient:
    if llm is None:
        raise RuntimeError(
            f"Missing LLM client for {reason}. Set OPENAI_API_KEY (and OPENAI_BASE_URL if needed), then retry."
        )
    return llm


def _table_to_markdown(table: list[list[str | None]]) -> str:
    clean_rows: list[list[str]] = []
    for row in table:
        clean = [(cell or "").replace("\n", " ").strip() for cell in row]
        if any(clean):
            clean_rows.append(clean)
    if not clean_rows:
        return ""

    width = max(len(row) for row in clean_rows)
    padded = [row + [""] * (width - len(row)) for row in clean_rows]
    header = padded[0]
    separator = ["---"] * width
    body = padded[1:] if len(padded) > 1 else []

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _extract_pdf(source_path: Path, *, llm: LLMClient | None, config: AppConfig) -> tuple[str, str]:
    try:
        import fitz
    except ImportError as exc:
        raise RuntimeError("Missing dependency: pymupdf. Run `pip install pymupdf`.") from exc
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("Missing dependency: pdfplumber. Run `pip install pdfplumber`.") from exc

    parts: list[str] = []
    low_text_pages: list[int] = []

    with pdfplumber.open(str(source_path)) as pdf:
        for page_index, page in enumerate(pdf.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            parts.append(f"## Page {page_index}")

            if page_text:
                parts.append(page_text)
            else:
                parts.append("(no direct text extracted)")

            tables = page.extract_tables() or []
            for table_index, table in enumerate(tables, start=1):
                as_markdown = _table_to_markdown(table)
                if as_markdown:
                    parts.append(f"\n### Table {table_index}\n{as_markdown}")

            if len(page_text) < 80:
                low_text_pages.append(page_index)

    if low_text_pages:
        llm_client = _require_llm(llm, reason="PDF OCR fallback")
        doc = fitz.open(str(source_path))
        temp_dir = Path(tempfile.mkdtemp(prefix="llm_wiki_pdf_"))
        try:
            for page_index in low_text_pages:
                page = doc[page_index - 1]
                pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
                image_path = temp_dir / f"page-{page_index}.png"
                pix.save(str(image_path))
                ocr_prompt = (
                    "Extract all readable text from this PDF page image and summarize layout.\n"
                    "Return markdown with:\n"
                    "## OCR Text\n"
                    "## Visual Notes\n"
                    "If unreadable, say so explicitly."
                )
                ocr_text = llm_client.describe_image(
                    model=config.vision_model,
                    image_path=image_path,
                    prompt=ocr_prompt,
                    max_tokens=min(1800, config.llm_max_tokens),
                )
                parts.append(f"\n## Page {page_index} (Vision OCR Fallback)\n{ocr_text}")
        finally:
            doc.close()
            shutil.rmtree(temp_dir, ignore_errors=True)

    return "\n\n".join(parts).strip(), "pdfplumber+llm_vision_fallback"


def _extract_docx(source_path: Path) -> tuple[str, str]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("Missing dependency: python-docx. Run `pip install python-docx`.") from exc

    doc = Document(str(source_path))
    parts: list[str] = []

    body_lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if body_lines:
        parts.append("## Body")
        parts.extend(f"- {line}" for line in body_lines)

    table_count = 0
    for table in doc.tables:
        table_count += 1
        rows: list[list[str]] = []
        for row in table.rows:
            rows.append([cell.text.strip().replace("\n", " ").strip() for cell in row.cells])
        table_md = _table_to_markdown(rows)
        if table_md:
            parts.append(f"\n## Table {table_count}\n{table_md}")

    header_idx = 0
    footer_idx = 0
    for section in doc.sections:
        header_text = "\n".join(p.text.strip() for p in section.header.paragraphs if p.text.strip()).strip()
        if header_text:
            header_idx += 1
            parts.append(f"\n## Header {header_idx}\n{header_text}")
        footer_text = "\n".join(p.text.strip() for p in section.footer.paragraphs if p.text.strip()).strip()
        if footer_text:
            footer_idx += 1
            parts.append(f"\n## Footer {footer_idx}\n{footer_text}")

    merged = "\n".join(parts).strip()
    if not merged:
        merged = "(no extractable text found in docx)"
    return merged, "python_docx_structured_extract"


def _extract_pptx(source_path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("Missing dependency: python-pptx. Run `pip install python-pptx`.") from exc

    prs = Presentation(str(source_path))
    parts: list[str] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {slide_index}")
        saw_content = False

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                rows: list[list[str]] = []
                for row in shape.table.rows:
                    rows.append([cell.text.strip().replace("\n", " ").strip() for cell in row.cells])
                table_md = _table_to_markdown(rows)
                if table_md:
                    parts.append(f"### Slide {slide_index} Table\n{table_md}")
                    saw_content = True
                continue

            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                text = shape.text_frame.text.strip()
            else:
                text = getattr(shape, "text", "")
                if isinstance(text, str):
                    text = text.strip()
                else:
                    text = ""

            if text:
                parts.append(f"- {text}")
                saw_content = True

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"- [Notes] {notes}")
                saw_content = True

        if not saw_content:
            parts.append("- (no text found)")

    return "\n".join(parts).strip(), "python_pptx_structured_extract"


def _extract_image(source_path: Path, *, llm: LLMClient | None, config: AppConfig) -> tuple[str, str]:
    prompt = (
        "You are extracting image knowledge for a compile-first RAG pipeline.\n"
        "Return markdown with sections:\n"
        "## Visual Summary\n"
        "## OCR Text\n"
        "## Objects And Entities\n"
        "## Facts / Claims Visible\n"
        "## Uncertainty\n"
        "Be faithful to the image. If text is unreadable, say so explicitly."
    )
    llm_client = _require_llm(llm, reason="image extraction")
    output = llm_client.describe_image(
        model=config.vision_model,
        image_path=source_path,
        prompt=prompt,
        max_tokens=min(2200, config.llm_max_tokens),
    )
    return output, "llm_vision_extract"


def _prepare_audio_for_asr(source_path: Path) -> tuple[Path, Path]:
    workdir = Path(tempfile.mkdtemp(prefix="llm_wiki_audio_"))
    normalized = workdir / "normalized.wav"
    cmd = [
        "ffmpeg",
        "-y",
        "-i",
        str(source_path),
        "-ac",
        "1",
        "-ar",
        "16000",
        "-vn",
        str(normalized),
    ]
    try:
        result = subprocess.run(cmd, capture_output=True, text=True)
    except FileNotFoundError as exc:
        raise RuntimeError("ffmpeg command not found. Please install FFmpeg and add it to PATH.") from exc
    if result.returncode != 0:
        raise RuntimeError(f"ffmpeg normalize failed: {result.stderr[:300]}")
    return normalized, workdir


def _extract_audio_or_video(source_path: Path, *, llm: LLMClient | None, config: AppConfig) -> tuple[str, str]:
    llm_client = _require_llm(llm, reason="audio/video transcription")
    normalized_path, workdir = _prepare_audio_for_asr(source_path)
    try:
        transcript = llm_client.transcribe_audio(model=config.asr_model, audio_path=normalized_path)
    finally:
        shutil.rmtree(workdir, ignore_errors=True)
    text = f"## Transcript\n{transcript}\n"
    return text, "ffmpeg_normalize+llm_asr"


def _extract_content(source_path: Path, *, llm: LLMClient | None, config: AppConfig) -> tuple[str, str, str]:
    ext = source_path.suffix.lower()

    if ext in TEXT_EXTENSIONS:
        body, extractor = _extract_text_like(source_path)
        return body, "text", extractor
    if ext in JSON_EXTENSIONS:
        body, extractor = _extract_json(source_path)
        return body, "json", extractor
    if ext in HTML_EXTENSIONS:
        body, extractor = _extract_html(source_path)
        return body, "html", extractor
    if ext in PDF_EXTENSIONS:
        body, extractor = _extract_pdf(source_path, llm=llm, config=config)
        return body, "pdf", extractor
    if ext in DOCX_EXTENSIONS:
        body, extractor = _extract_docx(source_path)
        return body, "docx", extractor
    if ext in PPTX_EXTENSIONS:
        body, extractor = _extract_pptx(source_path)
        return body, "pptx", extractor
    if ext in LEGACY_PPT_EXTENSIONS:
        raise RuntimeError("Legacy .ppt is not directly supported. Please convert to .pptx first.")
    if ext in IMAGE_EXTENSIONS:
        body, extractor = _extract_image(source_path, llm=llm, config=config)
        return body, "image", extractor
    if ext in AUDIO_EXTENSIONS or ext in VIDEO_EXTENSIONS:
        body, extractor = _extract_audio_or_video(source_path, llm=llm, config=config)
        media_type = "audio" if ext in AUDIO_EXTENSIONS else "video_audio"
        return body, media_type, extractor
    raise RuntimeError(f"Unsupported format: {ext}")


def ingest_sources(
    config: AppConfig,
    *,
    llm: LLMClient | None = None,
    inputs: list[str],
    force: bool = False,
    max_files: int | None = None,
    progress_cb: Callable[[int, int, str], None] | None = None,
) -> IngestSummary:
    ensure_workspace(config.raw_dir, config.wiki_dir, config.outputs_dir)

    all_files = resolve_sources(inputs, max_files=max_files)

    total = max(len(all_files), 1)
    if progress_cb:
        progress_cb(0, total, f"Resolved {len(all_files)} candidate files")

    ingested: list[IngestItem] = []
    skipped_existing: list[Path] = []
    skipped_unsupported: list[Path] = []
    errors: list[str] = []

    for idx, source in enumerate(all_files, start=1):
        ext = source.suffix.lower()
        if progress_cb:
            progress_cb(idx - 1, total, f"Ingesting {source.name}")
        if ext not in ALL_SUPPORTED_EXTENSIONS:
            skipped_unsupported.append(source)
            if progress_cb:
                progress_cb(idx, total, f"Unsupported format {source.suffix.lower()} for {source.name}")
            continue

        target = _target_path(config.raw_dir, source)
        if target.exists() and not force:
            skipped_existing.append(target)
            if progress_cb:
                progress_cb(idx, total, f"Skipped existing {target.name}")
            continue

        try:
            body, source_type, extractor = _extract_content(source, llm=llm, config=config)
            if not body.strip():
                raise RuntimeError("Extractor produced empty content.")

            content = (
                f"# Ingested Source: {source.name}\n\n"
                f"- source_path: `{source.as_posix()}`\n"
                f"- source_type: {source_type}\n"
                f"- extractor: {extractor}\n"
                f"- ingested_at: {now_stamp()}\n\n"
                f"## Extracted Content\n{body.strip()}\n"
            )
            write_text_file(target, content)
            ingested.append(
                IngestItem(
                    source_path=source,
                    output_path=target,
                    source_type=source_type,
                    extractor=extractor,
                )
            )
            if progress_cb:
                progress_cb(idx, total, f"Ingested {source.name} ({source_type})")
        except Exception as exc:
            errors.append(f"{source.as_posix()}: {exc}")
            if progress_cb:
                progress_cb(idx, total, f"Failed {source.name}: {exc}")

    if progress_cb:
        progress_cb(total, total, "Ingest stage finished")

    return IngestSummary(
        ingested=ingested,
        skipped_existing=skipped_existing,
        skipped_unsupported=skipped_unsupported,
        errors=errors,
    )
